from __future__ import annotations

from io import BytesIO
from pathlib import Path
from uuid import uuid4

from .job import utc_now
from .source_asset import SourceAssetType, SourceUnit


class SourceAssetParseError(Exception):
    pass


def parse_source_asset(
    asset_id: str,
    asset_type: SourceAssetType,
    content: bytes,
) -> tuple[list[SourceUnit], dict[str, object]]:
    if asset_type == "pptx":
        return _parse_pptx(asset_id, content)
    if asset_type == "pdf":
        return _parse_pdf(asset_id, content)
    if asset_type == "docx":
        return _parse_docx(asset_id, content)
    if asset_type == "text":
        return _parse_text(asset_id, content)
    raise SourceAssetParseError(f"Extraction is not supported for {asset_type}.")


def _parse_pptx(
    asset_id: str,
    content: bytes,
) -> tuple[list[SourceUnit], dict[str, object]]:
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise SourceAssetParseError("The PPTX file could not be parsed.") from exc

    units: list[SourceUnit] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        text = "\n".join(texts).strip()
        if text:
            units.append(
                _unit(
                    asset_id,
                    "slide",
                    len(units),
                    text,
                    {"slide_number": slide_number},
                )
            )
    return units, {"slide_count": len(presentation.slides)}


def _parse_pdf(
    asset_id: str,
    content: bytes,
) -> tuple[list[SourceUnit], dict[str, object]]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(content))
    except Exception as exc:
        raise SourceAssetParseError("The PDF file could not be parsed.") from exc

    units: list[SourceUnit] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        if text:
            units.append(
                _unit(
                    asset_id,
                    "page",
                    len(units),
                    text,
                    {"page_number": page_number},
                )
            )
    return units, {"page_count": len(reader.pages)}


def _parse_docx(
    asset_id: str,
    content: bytes,
) -> tuple[list[SourceUnit], dict[str, object]]:
    try:
        from docx import Document

        document = Document(BytesIO(content))
    except Exception as exc:
        raise SourceAssetParseError("The DOCX file could not be parsed.") from exc

    units = [
        _unit(
            asset_id,
            "paragraph",
            ordinal,
            paragraph.text.strip(),
            {"paragraph_number": ordinal + 1},
        )
        for ordinal, paragraph in enumerate(
            paragraph
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        )
    ]
    return units, {"paragraph_count": len(units)}


def _parse_text(
    asset_id: str,
    content: bytes,
) -> tuple[list[SourceUnit], dict[str, object]]:
    try:
        text = content.decode("utf-8-sig")
    except UnicodeDecodeError:
        try:
            text = content.decode("gb18030")
        except UnicodeDecodeError as exc:
            raise SourceAssetParseError(
                "The text file must use UTF-8 or GB18030 encoding."
            ) from exc
    chunks = _text_chunks(text)
    units = [
        _unit(asset_id, "text", index, chunk, {"section_number": index + 1})
        for index, chunk in enumerate(chunks)
    ]
    return units, {"section_count": len(units)}


def _text_chunks(text: str, *, max_characters: int = 4000) -> list[str]:
    paragraphs = [part.strip() for part in text.replace("\r\n", "\n").split("\n\n") if part.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        pieces = [
            paragraph[index:index + max_characters]
            for index in range(0, len(paragraph), max_characters)
        ] or [paragraph]
        for piece in pieces:
            candidate = f"{current}\n\n{piece}".strip() if current else piece
            if current and len(candidate) > max_characters:
                chunks.append(current)
                current = piece
            else:
                current = candidate
    if current:
        chunks.append(current)
    return chunks


def _unit(asset_id, unit_type, ordinal, text, locator) -> SourceUnit:
    return SourceUnit(
        id=uuid4().hex,
        asset_id=asset_id,
        unit_type=unit_type,
        ordinal=ordinal,
        text=text,
        locator=locator,
        created_at=utc_now(),
    )
