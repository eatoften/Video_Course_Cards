from io import BytesIO

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from app.source_asset_parser import parse_source_asset


def test_parse_pptx_preserves_slide_number():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    text_box.text = "Backpropagation applies the chain rule."
    content = BytesIO()
    presentation.save(content)

    units, metadata = parse_source_asset("asset-1", "pptx", content.getvalue())

    assert len(units) == 1
    assert units[0].unit_type == "slide"
    assert units[0].locator == {"slide_number": 1}
    assert "chain rule" in units[0].text
    assert metadata["slide_count"] == 1


def test_parse_docx_preserves_paragraph_number():
    document = Document()
    document.add_paragraph("Gradient descent updates model parameters.")
    content = BytesIO()
    document.save(content)

    units, metadata = parse_source_asset("asset-2", "docx", content.getvalue())

    assert len(units) == 1
    assert units[0].unit_type == "paragraph"
    assert units[0].locator == {"paragraph_number": 1}
    assert metadata["paragraph_count"] == 1


def test_parse_markdown_keeps_sections_local():
    units, metadata = parse_source_asset(
        "asset-3",
        "text",
        b"# Optimization\n\nMomentum smooths parameter updates.",
    )

    assert len(units) == 1
    assert "Momentum" in units[0].text
    assert metadata["section_count"] == 1
